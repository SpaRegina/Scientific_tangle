# -*- coding: utf-8 -*-
"""Извлечение текста из PDF/DOCX/PPTX/XLSX/TXT. Возвращает список страниц/блоков."""
import os
import re
import logging

log = logging.getLogger("kg.ingest")


def _clean(t: str) -> str:
    t = t.replace("\x00", " ")
    t = re.sub(r"-\n(?=[а-яa-z])", "", t)  # перенос слова
    t = re.sub(r"[ \t]+", " ", t)
    return t.strip()


def extract_pdf(path, max_pages=400):
    import fitz
    pages = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            pages.append((i + 1, _clean(page.get_text("text"))))
    return pages


def extract_docx(path):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    text = _clean("\n".join(parts))
    # разбивка на "страницы" по ~4000 символов
    return [(i + 1, text[i * 4000:(i + 1) * 4000]) for i in range((len(text) // 4000) + 1) if text[i * 4000:(i + 1) * 4000]]


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        pages.append((i + 1, _clean("\n".join(parts))))
    return pages


def extract_txt(path):
    with open(path, encoding="utf-8", errors="ignore") as f:
        text = _clean(f.read())
    return [(1, text)]


def _paginate(text, size=4000):
    return [(i + 1, text[i * size:(i + 1) * size])
            for i in range((len(text) // size) + 1) if text[i * size:(i + 1) * size]]


def extract_doc(path):
    """Старый .doc/.rtf — конвертация через LibreOffice (если установлен)."""
    import subprocess, tempfile, glob, shutil
    if not shutil.which("soffice"):
        return None
    with tempfile.TemporaryDirectory() as td:
        subprocess.run(["soffice", "--headless", "--convert-to", "txt:Text", "--outdir", td, path],
                       capture_output=True, timeout=180)
        fs = glob.glob(os.path.join(td, "*.txt"))
        if not fs:
            return None
        with open(fs[0], encoding="utf-8", errors="ignore") as f:
            return _paginate(_clean(f.read()))


def extract_xlsx(path):
    """Таблицы: каждая строка -> 'ячейка | ячейка'; лист = страница."""
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    pages = []
    for i, ws in enumerate(wb.worksheets[:20]):
        rows = []
        for row in ws.iter_rows(max_row=2000, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        pages.append((i + 1, _clean("\n".join(rows))))
    wb.close()
    return pages


def extract_xls(path):
    """Старый .xls — конвертация в .xlsx через LibreOffice, затем openpyxl."""
    import subprocess, tempfile, glob, shutil
    if not shutil.which("soffice"):
        return None
    with tempfile.TemporaryDirectory() as td:
        subprocess.run(["soffice", "--headless", "--convert-to", "xlsx", "--outdir", td, path],
                       capture_output=True, timeout=180)
        fs = glob.glob(os.path.join(td, "*.xlsx"))
        return extract_xlsx(fs[0]) if fs else None


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx, ".docm": extract_docx,
    ".pptx": extract_pptx,
    ".txt": extract_txt, ".md": extract_txt,
    ".doc": extract_doc, ".rtf": extract_doc,
    ".xlsx": extract_xlsx, ".xlsm": extract_xlsx,
    ".xls": extract_xls,
}
ARCHIVE_EXTS = {".zip", ".rar"}


def extract(path):
    """-> list[(page_no, text)] или None если формат не поддержан/ошибка."""
    ext = os.path.splitext(path)[1].lower()
    fn = EXTRACTORS.get(ext)
    if not fn:
        return None
    try:
        pages = fn(path)
        return [(n, t) for n, t in pages if t and len(t) > 30]
    except Exception as e:
        log.warning("Не удалось извлечь %s: %s", path, e)
        return None
